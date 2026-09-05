from __future__ import annotations

import io
import textwrap
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .extractors import VisualAsset
from .images import CommonsImage


SLIDE_W = 13.333
SLIDE_H = 7.5


def _safe_picture_bytes(data: bytes) -> bytes:
    with Image.open(io.BytesIO(data)) as img:
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGB")
        out = io.BytesIO()
        fmt = "PNG" if img.mode == "RGBA" else "JPEG"
        img.save(out, format=fmt, quality=92)
        return out.getvalue()


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.8), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 33, 45)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.78), Inches(1.45), Inches(11.5), Inches(0.5))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.color.rgb = RGBColor(80, 88, 99)


def _add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.72), Inches(7.1), Inches(11.9), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8.5)
    p.font.color.rgb = RGBColor(105, 110, 118)
    p.alignment = PP_ALIGN.RIGHT


def _add_notes(slide, notes: str) -> None:
    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is not None:
        notes_slide.notes_text_frame.text = notes.strip()


def _add_image_fit(slide, image_bytes: bytes, x: float, y: float, w: float, h: float) -> None:
    raw = _safe_picture_bytes(image_bytes)
    with Image.open(io.BytesIO(raw)) as img:
        iw, ih = img.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio > box_ratio:
        draw_w = w
        draw_h = w / image_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * image_ratio
        draw_y = y
        draw_x = x + (w - draw_w) / 2
    slide.shapes.add_picture(io.BytesIO(raw), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def build_pptx(
    plan: dict[str, Any],
    uploaded_images: list[VisualAsset],
    resolved_images: list[CommonsImage | None],
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    _add_title(title_slide, plan.get("deck_title", "Presentation"), plan.get("deck_subtitle", ""))
    tagline = title_slide.shapes.add_textbox(Inches(0.78), Inches(3.0), Inches(11.5), Inches(1.1))
    p = tagline.text_frame.paragraphs[0]
    p.text = "Õpetaja slaidid · automaatselt koostatud lähteallikate põhjal"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(70, 78, 90)
    _add_notes(title_slide, "Sissejuhatus: selgita teema eesmärki ja seosta see varasemate teadmistega.")

    for idx, slide_data in enumerate(plan.get("slides", [])):
        slide = prs.slides.add_slide(blank)
        _add_title(slide, str(slide_data.get("title", f"Slide {idx+1}")))

        bullets = [str(x) for x in slide_data.get("bullets", [])][:6]
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(7.25), Inches(4.95))
        tf = body.text_frame
        tf.word_wrap = True
        tf.clear()
        for b_idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(21)
            p.font.color.rgb = RGBColor(34, 42, 53)
            p.level = 0
            p.space_after = Pt(10)

        chosen_data: bytes | None = None
        credit = ""
        uploaded_index = slide_data.get("uploaded_image_index")
        try:
            ui = int(uploaded_index) - 1 if uploaded_index is not None else -1
        except (TypeError, ValueError):
            ui = -1
        if 0 <= ui < len(uploaded_images):
            chosen_data = uploaded_images[ui].data
            credit = f"Uploaded source: {uploaded_images[ui].name}"
        elif idx < len(resolved_images) and resolved_images[idx] is not None:
            web_img = resolved_images[idx]
            chosen_data = web_img.data
            credit = f"{web_img.credit} · {web_img.source_url}"

        if chosen_data:
            try:
                _add_image_fit(slide, chosen_data, 8.35, 1.72, 4.2, 4.8)
            except Exception:
                pass

        notes = str(slide_data.get("teacher_notes", ""))
        if credit:
            notes += f"\n\nIMAGE CREDIT: {credit}"
        _add_notes(slide, notes)
        _add_footer(slide, f"{plan.get('deck_title', '')} · {idx + 1}")

    refs = plan.get("references") or []
    if refs:
        slide = prs.slides.add_slide(blank)
        _add_title(slide, "Allikad ja lisalugemine")
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.1))
        tf = body.text_frame
        tf.clear()
        for idx, ref in enumerate(refs[:12]):
            title = str(ref.get("title") or ref.get("url") or "Source")
            url = str(ref.get("url") or "")
            used_for = str(ref.get("used_for") or "")
            line = f"{title} — {url}"
            if used_for:
                line += f" ({used_for})"
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = textwrap.shorten(line, width=180, placeholder="…")
            p.font.size = Pt(13)
            p.space_after = Pt(8)
        _add_notes(slide, "Kontrolli enne kasutamist eriti ajas muutuda võivaid fakte ja veebilinke.")

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
